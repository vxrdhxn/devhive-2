# pyre-ignore-all-errors
import io
import os
import csv
from typing import Optional
from backend.config import get_settings # type: ignore

settings = get_settings()

class DocumentProcessor:
    """Service for extracting text from various file formats"""
    
    @staticmethod
    def extract_text(file_content: bytes, filename: str) -> str:
        """Extract text based on file extension"""
        print(f"DEBUG: Extracting text from {filename}")
        extension = filename.split(".")[-1].lower()
        
        try:
            if extension == "pdf":
                return DocumentProcessor._extract_from_pdf(file_content)
            elif extension == "docx":
                return DocumentProcessor._extract_from_docx(file_content)
            elif extension == "pptx":
                return DocumentProcessor._extract_from_pptx(file_content)
            elif extension == "rtf":
                return DocumentProcessor._extract_from_rtf(file_content)
            elif extension == "xlsx":
                return DocumentProcessor._extract_from_xlsx(file_content)
            elif extension == "csv":
                return DocumentProcessor._extract_from_csv(file_content)
            elif extension in ["txt", "md"]:
                return DocumentProcessor._extract_from_text(file_content)
            else:
                raise ValueError(f"Unsupported file format: {extension}")
        except Exception as e:
            print(f"ERROR: Failed to extract text from {filename}: {str(e)}")
            raise ValueError(f"Failed to process {extension} file: {str(e)}")

    @staticmethod
    def _extract_from_pdf(file_content: bytes) -> str:
        """Extract text from PDF using PyPDF2"""
        import PyPDF2 # type: ignore
        pdf_file = io.BytesIO(file_content)
        reader = PyPDF2.PdfReader(pdf_file)
        text_parts = []
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)
        return "\n".join(text_parts).strip()

    @staticmethod
    def _extract_from_docx(file_content: bytes) -> str:
        """Extract text from DOCX using python-docx"""
        from docx import Document as DocxDocument # type: ignore
        docx_file = io.BytesIO(file_content)
        doc = DocxDocument(docx_file)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        return text.strip()

    @staticmethod
    def _extract_from_pptx(file_content: bytes) -> str:
        """Extract text from PPTX using python-pptx"""
        from pptx import Presentation # type: ignore
        pptx_file = io.BytesIO(file_content)
        presentation = Presentation(pptx_file)
        text_parts = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(getattr(shape, "text"))
        return "\n".join(text_parts).strip()

    @staticmethod
    def _extract_from_rtf(file_content: bytes) -> str:
        """Extract text from RTF using striprtf"""
        from striprtf.striprtf import rtf_to_text # type: ignore
        rtf_text = file_content.decode("utf-8")
        return rtf_to_text(rtf_text).strip()

    @staticmethod
    def _extract_from_xlsx(file_content: bytes) -> str:
        """Extract text from XLSX using openpyxl"""
        import openpyxl # type: ignore
        xlsx_file = io.BytesIO(file_content)
        wb = openpyxl.load_workbook(xlsx_file, data_only=True)
        text_parts = []
        for sheet in wb.worksheets:
            text_parts.append(f"Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join([str(cell) for cell in row if cell is not None])
                if row_text.strip():
                    text_parts.append(row_text)
        return "\n".join(text_parts).strip()

    @staticmethod
    def _extract_from_csv(file_content: bytes) -> str:
        """Extract text from CSV using csv module"""
        csv_text = file_content.decode("utf-8")
        reader = csv.reader(io.StringIO(csv_text))
        text_parts = []
        for row in reader:
            text_parts.append(" ".join(row))
        return "\n".join(text_parts).strip()

    @staticmethod
    def _extract_from_text(file_content: bytes) -> str:
        """Extract text from plain text or markdown files"""
        try:
            return file_content.decode("utf-8").strip()
        except UnicodeDecodeError:
            # Fallback to latin-1 if utf-8 fails
            return file_content.decode("latin-1").strip()

# Singleton instance
document_processor = DocumentProcessor()
